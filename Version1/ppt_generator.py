import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
import subprocess
import os
import streamlit as st
import ollama

# Function to interact with Ollama LLaMA 3.2 with clean formatting (no length limit for complete content)
def generate_with_llama(prompt):
    try:
        response = ollama.generate(model="llama3.2", prompt=f"{prompt} Provide *only* the concise, complete text or numbered list (no introductory phrases like 'Here is the output...', no asterisks, bold markers, or extra formatting). Ensure the response is detailed but fits a presentation slide, focusing on clarity and using bullet points for content slides.")
        text = response['response'].strip()
        return text
    except Exception as e:
        return f"Error with LLaMA: {str(e)}"

# Function to split text into 5 bullet points if needed
def split_into_bullets(text, max_points=5):
    lines = text.split('\n')
    if not lines or all(not line.strip() for line in lines):
        return ["No content available."]
    points = [line.strip() for line in lines if line.strip()]
    if len(points) > max_points:
        return points[:max_points]  # Limit to 5 points
    elif len(points) < max_points:
        # Pad with empty points or repeat if necessary
        while len(points) < max_points:
            points.append("Additional point to be expanded.")
        return points
    return points

# Function to ensure title has exactly 5 words
def ensure_title_5_words(title, col1, col2, user_prompt):
    words = title.split()
    if len(words) != 5:
        # If not exactly 5 words, generate a new title with exactly 5 words
        prompt = f"Generate a creative title for a presentation analyzing {col1} and {col2} based on this user expectation: '{user_prompt}'. Provide *only* a concise title with exactly 5 words, no formatting, no additional content."
        new_title = generate_with_llama(prompt)
        words = new_title.split()
        if len(words) != 5:
            return "Performance Insight"  # Fallback if LLaMA fails
        return new_title
    return title

# Backend EDA and slide generation function
def generate_eda_report(csv_file, col1, col2, plot_type, min_slides, user_prompt):
    # Reset file pointer and load CSV
    csv_file.seek(0)
    try:
        df = pd.read_csv(csv_file)
        if df.empty:
            return False, "CSV file is empty or has no valid data."
        num_rows = len(df)
        num_cols = len(df.columns)
    except Exception as e:
        return False, f"Error reading CSV: {str(e)}"

    # Initialize PowerPoint presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Helper functions
    def add_text_slide(prs, title, content, bg_color=None, use_bullets=True):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        if bg_color:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
        title_placeholder = slide.shapes.title
        title_placeholder.text = title[:50]  # Limit title length
        title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)  # Title font
        title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        title_placeholder.text_frame.paragraphs[0].font.bold = False  # Ensure normal text
        title_placeholder.height = Inches(1)  # Default title height
        title_placeholder.top = Inches(0)  # Position at top
        title_placeholder.left = Inches(1)  # Align left
        title_placeholder.width = Inches(8)  # Match content width
        body_placeholder = slide.placeholders[1]
        body_placeholder.width = Inches(8)
        body_placeholder.height = Inches(5)  # Default content height
        body_placeholder.left = Inches(1)
        body_placeholder.top = Inches(1)  # Position content immediately below title
        body_placeholder.text_frame.word_wrap = True  # Enable word wrap
        body_placeholder.text_frame.clear()  # Clear existing content
        if use_bullets:
            for point in content:
                if point.strip():
                    p = body_placeholder.text_frame.add_paragraph()
                    p.text = point.strip()
                    p.level = 0  # Bullet level
                    p.font.size = Pt(18)  # Body font
                    p.font.color.rgb = RGBColor(51, 51, 51)
                    p.font.bold = False  # Ensure normal text
                    p.space_after = Pt(10)
        else:
            body_placeholder.text = content  # Plain text for index
            for paragraph in body_placeholder.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(51, 51, 51)
                paragraph.font.bold = False
                paragraph.space_after = Pt(10)
        return slide

    def add_image_slide(prs, title, image_path, bg_color=None):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        if bg_color:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
        title_shape = slide.shapes.title
        title_shape.text = title[:50]  # Limit title length
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)  # Larger title font
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        title_shape.text_frame.paragraphs[0].font.bold = False  # Ensure normal text
        title_shape.top = Inches(0)  # Position at top
        title_shape.height = Inches(1)  # Default title height
        title_shape.left = Inches(1)  # Align left
        title_shape.width = Inches(8)  # Match content width
        left, top, width, height = Inches(1), Inches(1), Inches(8), Inches(5)  # Position image below title
        slide.shapes.add_picture(image_path, left, top, width, height)
        return slide

    # Generate content with LLaMA
    title_prompt = f"Generate a creative title for a presentation analyzing {col1} and {col2} based on this user expectation: '{user_prompt}'. Provide *only* the concise title with exactly 5 words, no formatting, no additional content."
    cover_title = ensure_title_5_words(generate_with_llama(title_prompt), col1, col2, user_prompt)

    # Create actual slide titles based on EDA and user prompt
    slide_titles = [
        "CSV Overview",  # Always included
        f"EDA Chart: {col1} vs {col2}",
        "EDA Statistics",
        "EDA Report",
        "Conclusion"
    ]
    current_slides = len(slide_titles)
    if current_slides < min_slides:
        for i in range(min_slides - current_slides):
            extra_prompt = f"Generate a concise title (under 50 characters, no formatting, no additional content) for an extra slide {i+1} analyzing {col1} and {col2}, based on this expectation: '{user_prompt}'."
            extra_title = generate_with_llama(extra_prompt)
            slide_titles.append(f"Extra Insight {i+1}: {extra_title}")

    # Generate index based on actual slide titles
    index_content = "\n".join(f"{i+1}. {title}" for i, title in enumerate(slide_titles))

    # Cover Slide (only title, centered vertically, exactly 5 words)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    title = slide.shapes.title
    title.text = cover_title
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title.text_frame.paragraphs[0].font.bold = False  # Ensure normal text
    title.height = Inches(1.5)  # Increased height for better centering
    title.top = Inches(3)  # Center vertically on 7.5-inch slide
    title.left = Inches(1)
    title.width = Inches(8)
    # Remove all other placeholders (e.g., subtitle) completely
    for shape in slide.shapes:
        if shape.placeholder_format.idx != 0:  # Keep only the title (placeholder 0)
            shape.element.getparent().remove(shape.element)

    # Index Slide (numbered list, no bullets)
    add_text_slide(prs, "Index", index_content, bg_color=RGBColor(245, 245, 220), use_bullets=False)

    # CSV Overview Slide
    csv_details = f"File: {os.path.basename(csv_file.name)}\nRows: {num_rows}\nCols: {num_cols}\nSelected: {col1}, {col2}"
    overview_points = [line.strip() for line in csv_details.split('\n') if line.strip()]
    add_text_slide(prs, slide_titles[0], overview_points, bg_color=RGBColor(245, 245, 220))

    # EDA of the Two Columns
    plt.figure(figsize=(8, 5))
    if pd.api.types.is_numeric_dtype(df[col1]) and pd.api.types.is_numeric_dtype(df[col2]):
        chart_path = "comparison_plot.png"
        if plot_type == "Scatter":
            df.plot.scatter(x=col1, y=col2, color='teal', alpha=0.5)
            plt.title(f"Scatter: {col1} vs {col2}", fontsize=12)
            plt.xlabel(col1, fontsize=10)
            plt.ylabel(col2, fontsize=10)
            plt.savefig(chart_path, bbox_inches='tight')
            plt.close()
            corr = df[[col1, col2]].corr().iloc[0, 1]
            stats1 = df[col1].describe()
            stats2 = df[col2].describe()
            stats_content = (
                f"{col1} Stats:\nMean: {stats1['mean']:.2f}, Std: {stats1['std']:.2f}\n"
                f"{col2} Stats:\nMean: {stats2['mean']:.2f}, Std: {stats2['std']:.2f}\n"
                f"Correlation: {corr:.2f}"
            )
        elif plot_type == "Hexbin":
            plt.hexbin(df[col1], df[col2], gridsize=20, cmap='Blues', mincnt=1)
            plt.colorbar(label='Count')
            plt.title(f"Hexbin: {col1} vs {col2}", fontsize=12)
            plt.xlabel(col1, fontsize=10)
            plt.ylabel(col2, fontsize=10)
            plt.savefig(chart_path, bbox_inches='tight')
            plt.close()
            corr = df[[col1, col2]].corr().iloc[0, 1]
            stats1 = df[col1].describe()
            stats2 = df[col2].describe()
            stats_content = (
                f"{col1} Stats:\nMean: {stats1['mean']:.2f}, Std: {stats1['std']:.2f}\n"
                f"{col2} Stats:\nMean: {stats2['mean']:.2f}, Std: {stats2['std']:.2f}\n"
                f"Correlation: {corr:.2f}"
            )
        elif plot_type == "Box":
            bins = pd.cut(df[col1], bins=3, precision=0)
            df_box = df[[col2]].copy()
            df_box['Binned_' + col1] = bins
            df_box.boxplot(column=col2, by='Binned_' + col1, grid=False, patch_artist=True)
            plt.title(f"Box: {col2} by Binned {col1}", fontsize=12)
            plt.suptitle('')
            plt.xlabel(f"Binned {col1}", fontsize=10)
            plt.ylabel(col2, fontsize=10)
            plt.xticks(rotation=45, ha='right', fontsize=8)
            plt.savefig(chart_path, bbox_inches='tight')
            plt.close()
            corr = df[[col1, col2]].corr().iloc[0, 1]
            stats1 = df[col1].describe()
            stats2 = df[col2].describe()
            stats_content = (
                f"{col1} Stats:\nMean: {stats1['mean']:.2f}, Std: {stats1['std']:.2f}\n"
                f"{col2} Stats:\nMean: {stats2['mean']:.2f}, Std: {stats2['std']:.2f}\n"
                f"Correlation: {corr:.2f}"
            )
        elif plot_type == "Bar":
            bins = pd.cut(df[col1], bins=3, precision=0)
            df.groupby(bins)[col2].mean().plot(kind='bar', color='lightcoral', edgecolor='black')
            plt.title(f"Bar: Mean {col2} by Binned {col1}", fontsize=12)
            plt.xlabel(f"Binned {col1}", fontsize=10)
            plt.ylabel(f"Mean {col2}", fontsize=10)
            plt.xticks(rotation=45, ha='right', fontsize=8)
            plt.savefig(chart_path, bbox_inches='tight')
            plt.close()
            corr = df[[col1, col2]].corr().iloc[0, 1]
            stats1 = df[col1].describe()
            stats2 = df[col2].describe()
            stats_content = (
                f"{col1} Stats:\nMean: {stats1['mean']:.2f}, Std: {stats1['std']:.2f}\n"
                f"{col2} Stats:\nMean: {stats2['mean']:.2f}, Std: {stats2['std']:.2f}\n"
                f"Correlation: {corr:.2f}"
            )
    else:
        if pd.api.types.is_numeric_dtype(df[col2]):
            df.groupby(col1)[col2].mean().plot(kind='bar', color='lightgreen', edgecolor='black')
            plt.title(f"Mean {col2} by {col1}", fontsize=12)
            plt.xlabel(col1, fontsize=10)
            plt.ylabel(f"Mean {col2}", fontsize=10)
            plt.xticks(rotation=45, ha='right', fontsize=8)
            chart_path = "comparison_bar.png"
            plt.savefig(chart_path, bbox_inches='tight')
            plt.close()
            stats_content = f"Mean {col2} by {col1}:\n{df.groupby(col1)[col2].mean().to_string()}"
        else:
            crosstab = pd.crosstab(df[col1], df[col2])
            crosstab.plot(kind='bar', stacked=True, colormap='Set2')
            plt.title(f"Cross-Tab: {col1} vs {col2}", fontsize=12)
            plt.xlabel(col1, fontsize=10)
            plt.ylabel("Count", fontsize=10)
            plt.xticks(rotation=45, ha='right', fontsize=8)
            chart_path = "comparison_crosstab.png"
            plt.savefig(chart_path, bbox_inches='tight')
            plt.close()
            stats_content = f"Cross-Tabulation:\n{crosstab.to_string()}"

    # Generate report/explanation with LLaMA (5 bullet points)
    report_prompt = f"Explain these statistics for {col1} and {col2} in 5 concise, complete bullet points (no introductory phrases or formatting): '{stats_content}'. Base it on this expectation: '{user_prompt}'."
    report_text = generate_with_llama(report_prompt)
    report_points = split_into_bullets(report_text, max_points=5)

    # Generate conclusion with LLaMA (5 bullet points)
    conclusion_prompt = f"Write a conclusion for a presentation analyzing {col1} and {col2} with at least {min_slides} slides, in 5 concise, complete bullet points (no introductory phrases or formatting), based on this expectation: '{user_prompt}' and these statistics: '{stats_content}'."
    conclusion_text = generate_with_llama(conclusion_prompt)
    conclusion_points = split_into_bullets(conclusion_text, max_points=5)

    # Add EDA slides with bullets
    add_image_slide(prs, slide_titles[1], chart_path, bg_color=RGBColor(240, 240, 240))
    add_text_slide(prs, slide_titles[2], stats_content.split('\n'), bg_color=RGBColor(240, 240, 240))
    add_text_slide(prs, slide_titles[3], report_points, bg_color=RGBColor(240, 240, 240))

    # Conclusion Slide
    add_text_slide(prs, slide_titles[4], conclusion_points, bg_color=RGBColor(240, 248, 255))

    # Add extra slides if needed
    for i in range(5, len(slide_titles)):
        extra_prompt = f"Generate content for slide '{slide_titles[i]}' analyzing {col1} and {col2}, in 5 concise, complete bullet points (no introductory phrases or formatting), based on this expectation: '{user_prompt}'."
        extra_text = generate_with_llama(extra_prompt)
        extra_points = split_into_bullets(extra_text, max_points=5)
        add_text_slide(prs, slide_titles[i], extra_points, bg_color=RGBColor(240, 240, 240))

    # Save and convert
    pptx_file = "report.pptx"
    odp_file = "report.odp"
    prs.save(pptx_file)
    try:
        subprocess.run(["libreoffice", "--headless", "--convert-to", "odp", pptx_file], check=True)
        if os.path.exists(chart_path):
            os.remove(chart_path)
        with open(odp_file, "rb") as f:
            odp_bytes = f.read()
        if os.path.exists(pptx_file):
            os.remove(pptx_file)
        if os.path.exists(odp_file):
            os.remove(odp_file)
        return True, odp_bytes
    except Exception as e:
        if os.path.exists(chart_path):
            os.remove(chart_path)
        return False, str(e)

# Streamlit front-end
def main():
    st.title("AI Based PPT Generator")
    st.markdown("Upload a CSV, set slide requirements, and provide a prompt for a custom EDA report in `.odp` format.")

    # File uploader
    uploaded_file = st.file_uploader("Choose a CSV file", type="csv")

    if uploaded_file is not None:
        uploaded_file.seek(0)
        try:
            df = pd.read_csv(uploaded_file)
            if df.empty:
                st.error("CSV file is empty or has no valid data.")
                return
            columns = list(df.columns)
            st.success(f"Loaded {uploaded_file.name} with {len(df)} rows and {len(columns)} columns.")
        except Exception as e:
            st.error(f"Error loading CSV: {str(e)}")
            return

        # User inputs
        col1 = st.selectbox("Select First Column", columns)
        col2 = st.selectbox("Select Second Column", columns)
        plot_type = st.selectbox("Select Plot Type (for numeric columns)", ["Scatter", "Hexbin", "Box", "Bar"])
        min_slides = st.number_input("Minimum Number of Slides", min_value=1, value=5, step=1)
        user_prompt = st.text_area("Describe how you want the PPT to be (e.g., 'professional report for stakeholders')", 
                                   "A concise and insightful analysis for a general audience", height=100)

        # Generate button
        if st.button("Generate Report"):
            if col1 == col2:
                st.warning("Please select two different columns.")
                return

            with st.spinner("Generating report with LLaMA..."):
                success, result = generate_eda_report(uploaded_file, col1, col2, plot_type, min_slides, user_prompt)
                if success:
                    st.success("Report generated successfully!")
                    st.download_button(
                        label="Download EDA Report",
                        data=result,
                        file_name="two_column_eda_report.odp",
                        mime="application/vnd.oasis.opendocument.presentation"
                    )
                else:
                    st.error(f"Error generating report: {result}")

if __name__ == "__main__":
    main()