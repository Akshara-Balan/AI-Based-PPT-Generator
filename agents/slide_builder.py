# agents/slide_builder.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
import random

class SlideBuilderAgent:
    def __init__(self, prs=None, theme="light"):
        self.prs = prs if prs else Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.theme = theme
        self.bg_colors = {
            "light": RGBColor(240, 240, 240),
            "dark": RGBColor(51, 51, 51),
            "blue": RGBColor(173, 216, 230),
            "green": RGBColor(144, 238, 144)
        }
        self.text_colors = {
            "light": RGBColor(51, 51, 51),
            "dark": RGBColor(255, 255, 255),
            "blue": RGBColor(0, 0, 139),
            "green": RGBColor(0, 100, 0)
        }
        self.title_colors = {
            "light": RGBColor(0, 51, 102),
            "dark": RGBColor(173, 216, 230),
            "blue": RGBColor(0, 0, 255),
            "green": RGBColor(0, 128, 0)
        }
        self.font_style = "Arial"  # Default

    def set_theme(self, theme):
        self.theme = theme

    def set_font_style(self, font_style):
        self.font_style = font_style

    def add_slide(self, title, content=None, chart_path=None, layout="text", table_data=None, progress=None):
        slide_layout = self.prs.slide_layouts[5] if chart_path else self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.bg_colors[self.theme]
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.title_colors[self.theme]
        title_shape.text_frame.paragraphs[0].font.name = self.font_style
        title_shape.top = Inches(0.5)
        title_shape.left = Inches(1)
        title_shape.width = Inches(8)
        title_shape.height = Inches(1)
        
        # Content based on layout
        if chart_path:
            slide.shapes.add_picture(chart_path, Inches(1), Inches(1.75), Inches(8), Inches(5))
        elif layout == "text" and content:
            textbox_height = Inches(6) if len(content) > 5 else Inches(5.5)
            font_size = Pt(14) if len(content) > 5 else Pt(16)
            tf = slide.shapes.add_textbox(Inches(1.5), Inches(1.75), Inches(7), textbox_height).text_frame
            tf.clear()
            for i, point in enumerate(content):
                p = tf.add_paragraph()
                p.text = f"• {point}"  # Icon-like bullet
                p.font.size = font_size
                p.font.color.rgb = self.text_colors[self.theme]
                p.font.name = self.font_style
                p.space_after = Pt(6) if len(content) > 5 else Pt(8)
        elif layout == "table" and table_data:
            rows, cols = len(table_data), len(table_data[0])
            table = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(1.75), Inches(7), Inches(4)).table
            for i, row in enumerate(table_data):
                for j, cell in enumerate(row):
                    table.cell(i, j).text = str(cell)
                    table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(14)
                    table.cell(i, j).text_frame.paragraphs[0].font.color.rgb = self.text_colors[self.theme]
                    table.cell(i, j).text_frame.paragraphs[0].font.name = self.font_style
        elif layout == "progress" and progress is not None:
            bar = slide.shapes.add_shape(1, Inches(1.5), Inches(2), Inches(7), Inches(0.5))  # Rectangle
            bar.fill.solid()
            bar.fill.fore_color.rgb = self.title_colors[self.theme]
            bar.width = Inches(7 * progress)
            text = slide.shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(7), Inches(0.5)).text_frame
            text.text = f"Progress: {progress*100:.0f}%"
            text.paragraphs[0].font.size = Pt(14)
            text.paragraphs[0].font.color.rgb = self.text_colors[self.theme]
            text.paragraphs[0].font.name = self.font_style
        
        return slide

    def add_title_slide(self, title, bg_color=RGBColor(240, 248, 255)):
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.title_colors[self.theme]
        title_shape.text_frame.paragraphs[0].font.name = self.font_style
        title_shape.top = Inches(3)
        title_shape.left = Inches(1)
        title_shape.width = Inches(8)
        for shape in slide.shapes:
            if shape.placeholder_format.idx != 0:
                shape.element.getparent().remove(shape.element)